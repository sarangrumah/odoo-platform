# -*- coding: utf-8 -*-
"""Plain-Python BRD text/structure extractor.

Not an Odoo model — kept as a normal class so it is unit-testable in isolation
and so the import errors for the optional native dependencies don't take down
the whole module load.

Each ``extract_*`` method returns a canonical dict:

    {
        "text": str,          # full markdown-ish concatenation, for fallback
        "sections": [
            {
                "title": str,
                "content": str,
                "level": int,    # 1 = H1, 2 = H2, ...
                "page": int,     # 1-based page or slide number
            },
            ...
        ],
    }
"""

from __future__ import annotations

import io
import logging
from typing import Any

_logger = logging.getLogger(__name__)


# Optional native deps — graceful degradation.
try:
    import fitz  # type: ignore  # PyMuPDF
except Exception:  # pragma: no cover
    fitz = None  # type: ignore

try:
    import docx  # type: ignore  # python-docx
except Exception:  # pragma: no cover
    docx = None  # type: ignore

try:
    from pptx import Presentation  # type: ignore  # python-pptx
except Exception:  # pragma: no cover
    Presentation = None  # type: ignore


class ExtractorDependencyError(RuntimeError):
    """Raised when the native dependency for the requested format is missing."""


class BrdExtractor:
    """Façade routing to the right format-specific extractor."""

    # ------------------------------------------------------------------
    # Dispatcher
    # ------------------------------------------------------------------

    def extract(self, binary: bytes, *, mime: str | None = None, filename: str | None = None) -> dict[str, Any]:
        kind = self._detect(mime=mime, filename=filename)
        if kind == "pdf":
            return self.extract_pdf(binary)
        if kind == "docx":
            return self.extract_docx(binary)
        if kind == "pptx":
            return self.extract_pptx(binary)
        raise ExtractorDependencyError(
            "Unsupported BRD file. Provide PDF, DOCX or PPTX. Got mime=%r filename=%r" % (mime, filename)
        )

    @staticmethod
    def _detect(*, mime: str | None, filename: str | None) -> str:
        m = (mime or "").lower()
        f = (filename or "").lower()
        if "pdf" in m or f.endswith(".pdf"):
            return "pdf"
        if "wordprocessingml" in m or f.endswith(".docx"):
            return "docx"
        if "presentationml" in m or f.endswith(".pptx"):
            return "pptx"
        return "unknown"

    # ------------------------------------------------------------------
    # PDF
    # ------------------------------------------------------------------

    def extract_pdf(self, binary: bytes) -> dict[str, Any]:
        if fitz is None:
            raise ExtractorDependencyError(
                "PyMuPDF is not installed in this Odoo runtime. Install with: pip install PyMuPDF"
            )
        sections: list[dict] = []
        full_text_parts: list[str] = []
        with fitz.open(stream=binary, filetype="pdf") as doc:  # type: ignore[union-attr]
            # First pass: use the embedded outline (TOC) if present, otherwise
            # fall back to "one section per page".
            toc = doc.get_toc() or []
            page_count = doc.page_count
            page_texts: list[str] = []
            for i in range(page_count):
                page_texts.append(doc.load_page(i).get_text("text"))
            full_text_parts.extend(page_texts)

            if toc:
                # toc entry: [level, title, page_number_1_based]
                for idx, entry in enumerate(toc):
                    level, title, page = entry[0], entry[1], entry[2]
                    end_page = toc[idx + 1][2] if idx + 1 < len(toc) else page_count
                    snippet = "\n".join(page_texts[max(0, page - 1) : max(0, end_page)])
                    sections.append(
                        {
                            "title": title.strip() or f"Section {idx + 1}",
                            "content": snippet.strip(),
                            "level": max(1, int(level)),
                            "page": int(page),
                        }
                    )
            else:
                for i, txt in enumerate(page_texts):
                    sections.append(
                        {
                            "title": f"Page {i + 1}",
                            "content": txt.strip(),
                            "level": 1,
                            "page": i + 1,
                        }
                    )
        return {"text": "\n\n".join(full_text_parts).strip(), "sections": sections}

    # ------------------------------------------------------------------
    # DOCX
    # ------------------------------------------------------------------

    def extract_docx(self, binary: bytes) -> dict[str, Any]:
        if docx is None:
            raise ExtractorDependencyError(
                "python-docx is not installed in this Odoo runtime. Install with: pip install python-docx"
            )
        document = docx.Document(io.BytesIO(binary))  # type: ignore[union-attr]
        sections: list[dict] = []
        full_text_parts: list[str] = []
        current: dict | None = None

        def _flush():
            nonlocal current
            if current is not None and (current["content"].strip() or current["title"].strip()):
                sections.append(current)
            current = None

        for para in document.paragraphs:
            text = para.text or ""
            full_text_parts.append(text)
            style = (para.style.name if para.style is not None else "") or ""
            if style.startswith("Heading"):
                _flush()
                try:
                    level = int(style.replace("Heading", "").strip() or "1")
                except ValueError:
                    level = 1
                current = {
                    "title": text.strip() or f"Heading {len(sections) + 1}",
                    "content": "",
                    "level": max(1, level),
                    "page": 1,
                }
            else:
                if current is None:
                    current = {"title": "Introduction", "content": "", "level": 1, "page": 1}
                if text.strip():
                    current["content"] += text + "\n"
        _flush()
        if not sections:
            sections.append({"title": "Document", "content": "\n".join(full_text_parts).strip(), "level": 1, "page": 1})
        return {"text": "\n".join(full_text_parts).strip(), "sections": sections}

    # ------------------------------------------------------------------
    # PPTX
    # ------------------------------------------------------------------

    def extract_pptx(self, binary: bytes) -> dict[str, Any]:
        if Presentation is None:
            raise ExtractorDependencyError(
                "python-pptx is not installed in this Odoo runtime. Install with: pip install python-pptx"
            )
        prs = Presentation(io.BytesIO(binary))  # type: ignore[misc]
        sections: list[dict] = []
        full_text_parts: list[str] = []
        for idx, slide in enumerate(prs.slides, start=1):
            title = ""
            body_parts: list[str] = []
            for shape in slide.shapes:
                if not getattr(shape, "has_text_frame", False):
                    continue
                # Title placeholder?
                is_title = False
                try:
                    if shape.is_placeholder and shape.placeholder_format is not None:
                        if shape.placeholder_format.idx == 0:
                            is_title = True
                except Exception:
                    is_title = False
                text = "\n".join(p.text for p in shape.text_frame.paragraphs if p.text)
                if is_title and not title:
                    title = text.strip()
                else:
                    if text.strip():
                        body_parts.append(text)
            content = "\n".join(body_parts).strip()
            sections.append(
                {
                    "title": title or f"Slide {idx}",
                    "content": content,
                    "level": 1,
                    "page": idx,
                }
            )
            full_text_parts.append(f"# {title or f'Slide {idx}'}\n{content}")
        return {"text": "\n\n".join(full_text_parts).strip(), "sections": sections}
