"""Build Erajaya VAS presentation (PPTX + PDF) with consistent branding.

PPTX  : python-pptx (handcrafted layouts, Erajaya red)
PDF   : reportlab (one slide per landscape A4 page)

Run:   python tools/build_presentation.py
Out :  docs/presentation-erajaya-vas.pptx
       docs/presentation-erajaya-vas.pdf
"""
from __future__ import annotations

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from reportlab.lib.pagesizes import landscape, A4
from reportlab.lib import colors as rl_colors
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import cm
from reportlab.platypus import (
    SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, PageBreak,
)
from reportlab.lib.enums import TA_LEFT, TA_CENTER

# Palette — neutral-dominant with red only as accent
BRAND = RGBColor(0xE3, 0x06, 0x13)            # Erajaya red — accents only
BRAND_DEEP = RGBColor(0xB8, 0x05, 0x0F)
BRAND_SOFT = RGBColor(0xFE, 0xE2, 0xE4)
SLATE = RGBColor(0x1F, 0x29, 0x37)            # primary chrome
SLATE_2 = RGBColor(0x37, 0x41, 0x51)          # secondary chrome
SLATE_3 = RGBColor(0x6B, 0x72, 0x80)          # tertiary
INK = RGBColor(0x11, 0x18, 0x27)
MUTED = RGBColor(0x6B, 0x72, 0x80)
SURFACE = RGBColor(0xF9, 0xFA, 0xFB)
SURFACE_2 = RGBColor(0xF3, 0xF4, 0xF6)
BORDER = RGBColor(0xE5, 0xE7, 0xEB)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

RL_BRAND = rl_colors.HexColor("#E30613")
RL_BRAND_DEEP = rl_colors.HexColor("#B8050F")
RL_BRAND_SOFT = rl_colors.HexColor("#FEE2E4")
RL_SLATE = rl_colors.HexColor("#1F2937")
RL_SLATE_2 = rl_colors.HexColor("#374151")
RL_SLATE_3 = rl_colors.HexColor("#6B7280")
RL_INK = rl_colors.HexColor("#111827")
RL_MUTED = rl_colors.HexColor("#6B7280")
RL_SURFACE = rl_colors.HexColor("#F9FAFB")
RL_SURFACE_2 = rl_colors.HexColor("#F3F4F6")
RL_BORDER = rl_colors.HexColor("#E5E7EB")

# ----------------------------------------------------------------------
# Slide content (single source — same data drives PPTX & PDF)
# ----------------------------------------------------------------------

SLIDES = []

def add(slide):
    SLIDES.append(slide)

# 1. Title
add({
    "kind": "title",
    "title": "Erajaya Value-Added Services",
    "subtitle": "Odoo Hub — Centralized ERP Platform for Erajaya Group",
    "footer": "Internal Business Presentation  ·  Product Owner VAS  ·  Mei 2026",
})

# 2. Executive Summary
add({
    "kind": "bullets",
    "title": "Executive Summary",
    "intro": "Odoo Hub: platform terpusat untuk percepatan delivery ERP lintas vertikal Erajaya Group.",
    "bullets": [
        "Centralized Module Repository — 1 sumber kebenaran modul, reusable lintas vertikal",
        "Simplified Deployment — provisioning tenant baru otomatis via orchestrator",
        "Centralized Monitoring — 1 control plane untuk seluruh tenant operasional",
        "Mandays implementasi turun signifikan karena modul standar siap pakai",
    ],
    "highlights": [
        ("82", "Modul Siap Pakai"),
        ("~70%", "Modul Reusable"),
        ("~77", "Mandays / Tenant"),
        ("99.5%", "SLA Target"),
    ],
})

# 3. Hub Platform Stack
add({
    "kind": "layers",
    "title": "Hub Platform Stack",
    "layers": [
        ("AI Layer",
         "ai-gateway (Claude / OpenAI / Ollama) · Ask-AI, anomaly inbox, NLQ, OCR receipt, churn prediction"),
        ("Observability Plane",
         "Prometheus · Grafana · Loki · Alertmanager · custom-predictor (capacity forecast lintas tenant)"),
        ("Multi-Tenant Runtime",
         "Odoo 19 + 82 custom modules · DB-per-tenant · Caddy/nginx + TLS · Redis · object storage"),
        ("Tenant Orchestrator",
         "FastAPI service: SSH bootstrap, Docker stack, DB create, module install, mail config — semua otomatis"),
        ("Centralized Module Repository",
         "Single source of truth 82 modul (core · ee-gap · localized · ops · vertical template) · LGPL-3 · CI-tested"),
    ],
})

# 4. High-Level Architecture (Linux-based)
add({
    "kind": "diagram",
    "title": "High-Level Architecture (Linux-based)",
    "intro": "Stack berjalan di atas Linux host dengan Docker — open, portable, dan tanpa lock-in OS.",
    "ascii": [
        "┌──────────────────────────────────────────────────────┐",
        "│  Linux Host  (Ubuntu 22.04 LTS · bare-metal / VPS)   │",
        "├──────────────────────────────────────────────────────┤",
        "│  Docker Engine 24+  ·  Compose v2  ·  systemd        │",
        "├──────────────────────────────────────────────────────┤",
        "│  Container Network (bridge)                          │",
        "│                                                      │",
        "│  ┌────────┐   ┌────────┐   ┌────────┐                │",
        "│  │ Caddy  │ → │  Odoo  │ ← │ AI GW  │                │",
        "│  │ TLS LB │   │workers │   │FastAPI │                │",
        "│  └───┬────┘   └───┬────┘   └───┬────┘                │",
        "│      │            │            │                     │",
        "│  ┌───▼────┐   ┌───▼────┐   ┌───▼────┐                │",
        "│  │ Redis  │   │Postgres│   │ Ollama │                │",
        "│  │ cache  │   │  15+   │   │ local  │                │",
        "│  └────────┘   └────────┘   └────────┘                │",
        "│                                                      │",
        "│  ┌──────────┐ ┌─────────┐ ┌──────┐ ┌─────────────┐   │",
        "│  │Prometheus│ │ Grafana │ │ Loki │ │Alertmanager │   │",
        "│  └──────────┘ └─────────┘ └──────┘ └─────────────┘   │",
        "├──────────────────────────────────────────────────────┤",
        "│  Persistent Volumes  (filestore · DB · logs · backup)│",
        "├──────────────────────────────────────────────────────┤",
        "│  Kernel hardening: AppArmor · seccomp · namespaces   │",
        "└──────────────────────────────────────────────────────┘",
    ],
    "decisions": [
        "OS: Ubuntu 22.04 LTS — 5 tahun security update",
        "Runtime: Docker 24+ + Compose v2, non-root user di tiap container",
        "Image hardening: distroless / Alpine base, CIS-aligned",
        "Storage: persistent volumes (filestore, DB, log, backup) bind-mount",
        "Network: bridge internal; hanya Caddy expose port 80/443",
        "Firewall: UFW / iptables — port 22 (SSH ops VPN), 80, 443 saja",
        "TLS: Caddy ACME (Let's Encrypt) auto-renew",
        "Patch: unattended-upgrades untuk security patches OS",
        "Backup: pg_dumpall nightly + filestore rsync ke object storage",
        "Portable: stack sama untuk on-prem, VPS, atau cloud (AWS/GCP/Azure)",
    ],
})

# 5. Centralized Module Repository — Reuse Matrix
add({
    "kind": "table",
    "title": "Centralized Module Repository — Reuse Matrix",
    "intro": "Modul kunci dipakai ulang lintas 6 vertikal Erajaya — bukan develop ulang per tenant.",
    "headers": ["Module", "F&B", "Active", "Eraspace", "Distrib", "Service", "Corp"],
    "rows": [
        ["custom_core",             "✓", "✓", "✓", "✓", "✓", "✓"],
        ["custom_accounting_full",  "✓", "✓", "✓", "✓", "✓", "✓"],
        ["custom_hr_payroll_id",    "✓", "✓", "✓", "✓", "✓", "✓"],
        ["custom_attendance",       "✓", "✓", "✓", "✓", "✓", "✓"],
        ["custom_approval_engine",  "✓", "✓", "✓", "✓", "✓", "✓"],
        ["custom_pdp_* (6 modul)",  "✓", "✓", "✓", "✓", "✓", "✓"],
        ["custom_coretax + bupot",  "✓", "✓", "✓", "✓", "✓", "✓"],
        ["custom_helpdesk",         "✓", "✓", "✓", "✓", "✓", "✓"],
        ["custom_whatsapp",         "✓", "✓", "✓", "✓", "✓", "✓"],
        ["custom_pos_id",           "✓", "✓", "✓", "—", "—", "—"],
        ["custom_ecommerce",        "✓", "✓", "✓", "—", "—", "—"],
        ["custom_wms_* (3 modul)",  "—", "✓", "✓", "✓", "✓", "—"],
        ["custom_field_service",    "—", "—", "—", "—", "✓", "—"],
        ["custom_subscription",     "—", "—", "—", "—", "—", "✓"],
    ],
    "col_widths": [3.4, 0.95, 0.95, 1.0, 0.95, 0.95, 0.95],
    "footnote": "~70% modul shared lintas vertikal · ~30% extension vertical-specific. Modul baru ditambahkan sekali di repo → langsung tersedia untuk semua tenant.",
})

# 5. Module Library — Capability Highlights
add({
    "kind": "modules_p2",
    "title": "Module Library — Capability Highlights",
    "groups": [
        ("Finance & Tax (Indonesian-localized, ready)",
         "PSAK 5-digit CoA · Intercompany & consolidation · Fixed asset depreciation · PPh 21 TER · BPJS Kes/TK · SPT 1721 A1 · PPN DPP Nilai Lain · e-Faktur Coretax · Bupot Unifikasi — semua built-in, tidak perlu develop per vertikal"),
        ("Human Capital (Indonesian-localized, ready)",
         "Geofence attendance · Cuti UU Cipta Kerja · Performance appraisal 360 · Recruitment + job-board webhook · Expense OCR receipt · Billable timesheet -> payroll"),
        ("Sales · CRM · Commerce",
         "Predictive lead scoring · Drip campaigns · Mass email · Midtrans/Xendit/DOKU · JNE/JNT/SiCepat/AnterAja · Subscription MRR · Asset rental w/ BAST · WhatsApp QR event ticket"),
        ("Service Operations",
         "Helpdesk SLA + escalation · Field Service dispatch · Repairs w/ warranty · Appointments · Livechat -> AI suggested reply · Frontdesk visitor"),
        ("Manufacturing & WMS",
         "MRP PLM (ECO + BoM versioning) · Quality + CAPA · Maintenance MTBF/MTTR · WMS putaway/cycle-count/to-engine · Mobile barcode · Zebra HHT PWA · IoT webhook"),
        ("Productivity & Cross-Cutting",
         "Studio-Lite (no-code customization) · Dashboards KPI + AI NLQ · Spreadsheet · Documents · E-signature · Knowledge wiki · Generic approval engine (delegation + OOO + SLA escalation)"),
    ],
})

# 6. Indonesian Localization Ready (merge dari 2 compliance slide)
add({
    "kind": "table",
    "title": "Indonesian Localization Ready",
    "intro": "Aturan akunting, perpajakan, ketenagakerjaan, dan data protection Indonesia — built-in, package siap pakai per tenant.",
    "headers": ["Domain", "Cakupan Localized"],
    "rows": [
        ["Akunting (PSAK)",
         "Chart of Accounts 5-digit aligned PSAK · Intercompany automation · Consolidation + eliminations · Fixed asset depreciation · Faktur Pengganti workflow"],
        ["Perpajakan (DJP)",
         "e-Faktur Coretax (NSFP 17 digit PER-11/PJ/2025) · Bupot PPh 21/23/26/Unifikasi · PPh 21 TER (PP 58/2023) · PPN DPP Nilai Lain (PMK 131/2024) · Sertel Fernet-encrypted · Pajakku ASPP H2H adapter"],
        ["HR & Ketenagakerjaan",
         "BPJS Kesehatan & Ketenagakerjaan (JHT/JKK/JKM/JP) · PTKP & THR · SPT 1721 A1 · Cuti UU Cipta Kerja · Payslip approval flow"],
        ["Data Protection (UU PDP)",
         "Klasifikasi data field-level · Consent management · DSAR endpoint · Audit log append-only hash-chained · PII masking · Retention auto-purge"],
    ],
    "col_widths": [2.5, 6.5],
    "footnote": "Pendekatan adapter pattern: ASP/regulasi berubah → swap implementasi tanpa ubah workflow tenant.",
})

# 7. Simplified Deployment — Tenant Orchestrator
add({
    "kind": "diagram",
    "title": "Simplified Deployment — Tenant Orchestrator",
    "intro": "Provisioning tenant baru otomatis, estimasi ~30 menit dari permintaan ke siap UAT.",
    "ascii": [
        "  Hub-Portal  (Vite + React 18 · port 18000)",
        "  Operator klik 'Provision Tenant'",
        "                       │",
        "  Tenant Orchestrator (FastAPI)",
        "                       │",
        "          ┌────────────┼────────────┐",
        "          ▼            ▼            ▼",
        "      tenant1 DB   tenant2 DB   tenantN DB",
        "      (isolated)   (isolated)   (isolated)",
        "",
        "  Shared plane:",
        "    • ai-gateway   • Postgres cluster   • Redis",
        "    • Object storage (filestore, sertel)",
        "    • Observability (Prom/Grafana/Loki)",
    ],
    "decisions": [
        "Step 1 — SSH bootstrap target VPS (Docker + Caddy install)",
        "Step 2 — Pull stack: Odoo + Postgres + Redis + module repo",
        "Step 3 — Create database tenant + apply addons path",
        "Step 4 — Install modul standar (sesuai profile vertikal)",
        "Step 5 — Generate Caddy route + TLS otomatis (ACME)",
        "Step 6 — Konfigurasi mail (SMTP / IMAP) + integrasi (Pajakku, payment)",
        "Step 7 — Smoke test + handover ke PO untuk UAT",
        "Total: ~30 menit / tenant — tidak perlu DevOps mandays manual",
    ],
})

# 8. Centralized Monitoring (NEW)
add({
    "kind": "two_col",
    "title": "Centralized Monitoring",
    "left_title": "Apa yang Dimonitor",
    "left_bullets": [
        "Odoo runtime metrics per tenant (requests, latency, error rate)",
        "Database health (connections, slow query, replication lag)",
        "AI gateway cost & latency per tenant",
        "Pajakku ASPP circuit state (open / closed / half-open)",
        "Audit chain integrity (PDP hash-chain verifier nightly)",
        "Tenant resource usage (CPU, memory, disk, filestore size)",
        "Capacity forecast 7-hari via custom-predictor",
        "TLS expiry, backup status, scheduler health",
    ],
    "right_title": "Manfaat Operasional",
    "right_bullets": [
        "1 dashboard Grafana untuk N tenant — bukan login per server",
        "Alert centralized via Alertmanager → ops on-call",
        "Proactive scaling — predictor rekomendasi upgrade hardware sebelum bottleneck",
        "MTTR turun — runbook + log + metric satu tempat",
        "Tenant SLA visible — laporan uptime per bulan otomatis",
        "Cost attribution per tenant (AI, storage, compute)",
        "Audit-ready: alert log + immutable trail",
    ],
    "footnote": "Stack: Prometheus (scrape 15s) · Grafana · Loki · Alertmanager · custom-predictor sidecar.",
})

# 9. Security Posture (slim)
add({
    "kind": "two_col",
    "title": "Security Posture",
    "left_title": "Application & Data Isolation",
    "left_bullets": [
        "DB-per-tenant isolation — bukan schema-per-tenant",
        "RBAC: Odoo groups + record rules per modul",
        "Multi-tier approval w/ delegation, OOO, SLA escalation",
        "Append-only audit log + PostgreSQL trigger",
        "Tenant allow-list per request (HMAC-validated)",
        "Secrets via SOPS-encrypted di repo + Fernet for sertel",
    ],
    "right_title": "Infrastructure & Pipeline",
    "right_bullets": [
        "CIS-hardened distroless containers, non-root",
        "AppArmor + seccomp profiles",
        "TLS termination + HSTS (Caddy / nginx)",
        "Pre-commit: gitleaks, ruff, bandit, hadolint",
        "CI: Semgrep (SAST) + pip-audit + Trivy + cosign signing",
        "Nightly pg_dumpall + DR runbook (drill executed Q2 2026)",
    ],
})

# 10. AI Layer
add({
    "kind": "two_col",
    "title": "AI Layer — Tertanam di Seluruh ERP",
    "left_title": "Infrastruktur AI",
    "left_bullets": [
        "ai-gateway (FastAPI sidecar) — multi-provider abstraction",
        "Provider switch via env: Claude / OpenAI / Ollama (local)",
        "HMAC-validated Odoo -> gateway calls",
        "Prompt caching untuk efisiensi cost",
        "Per-tenant rate limit & quota",
        "custom-predictor — tabular ML, capacity forecasting 7-hari",
    ],
    "right_title": "Fitur AI di Modul Bisnis",
    "right_bullets": [
        "Ask AI server-action di 9 model utama (invoice, payslip, picking)",
        "Anomaly Inbox — scan harian + severity + suggested action",
        "NLQ Chat — query natural language dengan PDP masking",
        "Document Auto-Classify — saran tag + classification",
        "AI churn prediction (subscription)",
        "AI suggested reply (livechat), moderation (forum)",
        "AI OCR receipt (expense), task breakdown (todo)",
        "Predictive lead scoring (CRM), spreadsheet helpers",
    ],
})

# 11. Business Value — Implementation Mandays (NEW)
add({
    "kind": "table",
    "title": "Business Value — Implementation Mandays per Tenant",
    "intro": "Estimasi standar onboarding 1 vertikal baru di Odoo Hub dengan semua modul standar siap pakai. Angka dalam tanda kurung = mandays opsional (hanya jika ada custom development).",
    "headers": ["Implementation Phase", "ITBP", "PO", "DevOps", "Developer (opt)"],
    "rows": [
        ["Requirement gathering / BRD",                 "5",  "8",  "—",  "—"],
        ["Tenant provisioning (orchestrator)",          "—",  "1",  "2",  "—"],
        ["Module activation (standard)",                "3",  "5",  "2",  "—"],
        ["Master data setup (CoA, partner, employee)",  "5",  "3",  "—",  "—"],
        ["Configuration (RBAC, approval, workflow)",    "3",  "5",  "1",  "—"],
        ["Integration setup (Pajakku, payment, WA)",    "2",  "2",  "3",  "(2)"],
        ["Custom module development",                   "—",  "(3)","—",  "(15)"],
        ["UAT + bug-fix",                               "5",  "5",  "1",  "(5)"],
        ["Training",                                    "3",  "5",  "—",  "—"],
        ["Go-live + hypercare",                         "3",  "3",  "2",  "—"],
        ["Subtotal — tanpa custom dev",                 "29", "37", "11", "0"],
        ["Subtotal — dengan custom dev",                "29", "40", "11", "22"],
    ],
    "col_widths": [4.0, 1.0, 1.0, 1.0, 1.4],
    "footnote": "Standar: ~77 mandays/tenant (lintas 3 role). Dengan 1 custom module: ~102 mandays. Baseline implementasi Odoo tradisional tanpa hub: ~200-300 mandays. Saving 60-70% per onboarding setelah library modul ready.",
})

# 12. VAS Productization
add({
    "kind": "table",
    "title": "VAS Productization — 5 Product Lines",
    "intro": "Platform Hub membuka monetisasi ke ekosistem Erajaya & eksternal.",
    "headers": ["VAS Product", "Modul Penopang", "Target Market"],
    "rows": [
        ["ERP-as-a-Service Multi-Tenant", "Hub stack lengkap + orchestrator", "Internal tenants & afiliasi Erajaya"],
        ["Localized Compliance Bundle", "custom_coretax* + Pajakku + custom_pdp_* + custom_hr_payroll_id", "UMKM ekosistem Erajaya (Eraspace partners, F&B franchisee)"],
        ["AI Operations Layer", "ai-gateway + custom_ai_features + predictor", "Bundled tiap tenant; advanced tier opt-in"],
        ["HHT / Field Ops Bridge", "custom_hht_bridge (Zebra PWA) + field_service", "Service Center, distribution, warehouse"],
        ["Vertical Template Accelerator", "Repo modul + onboarding journey", "Vertikal baru di dalam Erajaya Group"],
    ],
    "col_widths": [2.6, 3.4, 3.0],
})

# 13. Roadmap
add({
    "kind": "roadmap",
    "title": "Roadmap",
    "quarters": [
        ("Q2 2026 — DONE / IN PROGRESS", [
            "82 modul base shipping",
            "Multi-tenant orchestrator",
            "Pajakku ASPP adapter",
            "Hub-Portal UI (Vite+React)",
            "Production drill #1 executed",
        ]),
        ("Q3 2026", [
            "Onboard 3 internal vertical (F&B, Service Center, Corp Svcs)",
            "Centralized monitoring dashboard rollout",
            "HHT rollout untuk service vertical",
            "AI cost optimization (cache hit target 60%)",
        ]),
        ("Q4 2026", [
            "Onboard tenant 4–10 (target Y1)",
            "Module reuse audit + dedup",
            "ESG report otomatis (POJK 51/2017)",
            "Marketplace add-on (3rd-party vertical modules)",
        ]),
        ("2027", [
            "Tenant 11–25 (post-ROI threshold)",
            "Open partner program (ASPP, integrator)",
            "Localized Compliance Bundle GTM eksternal",
        ]),
    ],
})

# 14. Risks (slim, 5 paling material)
add({
    "kind": "table",
    "title": "Risks & Mitigations",
    "headers": ["Risk", "Mitigation"],
    "rows": [
        ["Odoo 19 upstream upgrade breaks customs",
         "custom_studio_lite declarative, modules pinned 19.0.x.y, CI per release"],
        ["Pajakku API / Coretax submission gagal",
         "Circuit breaker + manual fallback Coretax portal + ops alert"],
        ["Tenant data leakage",
         "DB-per-tenant + HMAC + tenant allow-list di ai-gateway"],
        ["AI cost explosion",
         "Per-tenant quota + prompt cache + local Ollama fallback"],
        ["Mandays meleset dari estimasi karena scope creep",
         "BRD freeze sebelum provisioning · change request via approval engine · re-baseline per onboarding"],
    ],
    "col_widths": [3.5, 5.5],
})

# 15. Ask / CTA
add({
    "kind": "cta",
    "title": "Call to Action",
    "intro": "Untuk Decision Maker Erajaya — endorsement & pilot funding.",
    "items": [
        ("1", "Endorsement",
         "Odoo Hub sebagai backbone delivery ERP Erajaya Group"),
        ("2", "Pilot Funding",
         "3 vertikal pertama (F&B, Service Center, Corporate Services) — Q3 2026"),
        ("3", "Mandays Baseline Calibration",
         "Lock estimasi mandays per vertikal selama 2 onboarding pertama untuk kalibrasi"),
        ("4", "Hiring Runway",
         "2 senior backend, 1 SRE, 1 PO untuk scale Y1 -> Y2"),
        ("5", "GTM Eksternal",
         "Localized Compliance Bundle sebagai produk eksternal Q4 2026"),
    ],
})

# 16. Thank you
add({
    "kind": "closing",
    "title": "Terima Kasih",
    "subtitle": "Diskusi & Q&A",
    "footer": "Product Owner — Value-Added Services  ·  Erajaya  ·  Mei 2026",
})


# ======================================================================
# PPTX builder
# ======================================================================

DEFAULT_FOOTER = "Erajaya VAS  ·  Odoo 19 Platform  ·  Mei 2026"


def build_pptx(out_path: Path, slides=None, footer_text: str = DEFAULT_FOOTER):
    if slides is None:
        slides = SLIDES
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    SW, SH = prs.slide_width, prs.slide_height

    blank = prs.slide_layouts[6]

    def add_rect(slide, left, top, width, height, fill, line=None):
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
        if line is None:
            shape.line.fill.background()
        else:
            shape.line.color.rgb = line
        shape.shadow.inherit = False
        return shape

    def add_text(slide, left, top, width, height, text, *,
                 size=14, bold=False, color=INK, align=PP_ALIGN.LEFT,
                 anchor=MSO_ANCHOR.TOP, font="Segoe UI"):
        tb = slide.shapes.add_textbox(left, top, width, height)
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = Emu(0)
        tf.margin_right = Emu(0)
        tf.margin_top = Emu(0)
        tf.margin_bottom = Emu(0)
        tf.vertical_anchor = anchor
        lines = text.split("\n") if isinstance(text, str) else text
        for i, line in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = align
            run = p.add_run()
            run.text = line
            run.font.name = font
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = color
        return tb

    def header_band(slide, title, subtitle=None):
        # Top band — slate with thin red accent line
        add_rect(slide, 0, 0, SW, Inches(0.9), SLATE)
        add_rect(slide, 0, Inches(0.9), SW, Inches(0.04), BRAND)
        add_text(slide, Inches(0.5), Inches(0.18), SW - Inches(1), Inches(0.55),
                 title, size=24, bold=True, color=WHITE)
        # Footer band
        add_rect(slide, 0, SH - Inches(0.35), SW, Inches(0.35), SURFACE_2)
        add_text(slide, Inches(0.5), SH - Inches(0.31), SW - Inches(1), Inches(0.28),
                 footer_text,
                 size=9, color=MUTED, align=PP_ALIGN.LEFT)
        add_text(slide, Inches(0.5), SH - Inches(0.31), SW - Inches(1), Inches(0.28),
                 "Confidential — Internal", size=9, color=MUTED, align=PP_ALIGN.RIGHT)

    def build_table(slide, left, top, width, height, headers, rows,
                    col_widths=None, header_fill=SLATE, header_fg=WHITE,
                    row_fill_alt=SURFACE_2, font_size=10):
        n_cols = len(headers)
        n_rows = len(rows) + 1
        tbl_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
        tbl = tbl_shape.table
        # Column widths
        if col_widths:
            total = sum(col_widths)
            for i, w in enumerate(col_widths):
                tbl.columns[i].width = int(width * w / total)
        # Headers
        for j, h in enumerate(headers):
            cell = tbl.cell(0, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = header_fill
            cell.text = ""
            tf = cell.text_frame
            tf.margin_left = Inches(0.08)
            tf.margin_right = Inches(0.08)
            tf.margin_top = Inches(0.04)
            tf.margin_bottom = Inches(0.04)
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = h
            run.font.name = "Segoe UI"
            run.font.size = Pt(font_size + 1)
            run.font.bold = True
            run.font.color.rgb = header_fg
        # Body
        for ri, row in enumerate(rows, start=1):
            for j, val in enumerate(row):
                cell = tbl.cell(ri, j)
                cell.fill.solid()
                cell.fill.fore_color.rgb = SURFACE_2 if ri % 2 == 1 else WHITE
                cell.text = ""
                tf = cell.text_frame
                tf.margin_left = Inches(0.08)
                tf.margin_right = Inches(0.08)
                tf.margin_top = Inches(0.04)
                tf.margin_bottom = Inches(0.04)
                tf.word_wrap = True
                p = tf.paragraphs[0]
                run = p.add_run()
                run.text = str(val)
                run.font.name = "Segoe UI"
                run.font.size = Pt(font_size)
                run.font.color.rgb = INK
        return tbl

    for s in slides:
        slide = prs.slides.add_slide(blank)
        kind = s["kind"]

        if kind == "title":
            add_rect(slide, 0, 0, SW, SH, WHITE)
            # Slate left band + thin red accent stripe
            add_rect(slide, 0, 0, Inches(0.5), SH, SLATE)
            add_rect(slide, Inches(0.5), 0, Inches(0.06), SH, BRAND)
            add_text(slide, Inches(1.2), Inches(2.4), SW - Inches(2.4), Inches(1.2),
                     s["title"], size=44, bold=True, color=INK)
            add_text(slide, Inches(1.2), Inches(3.6), SW - Inches(2.4), Inches(0.8),
                     s["subtitle"], size=24, color=SLATE_2)
            add_rect(slide, Inches(1.2), Inches(4.5), Inches(1.8), Inches(0.05), BRAND)
            add_text(slide, Inches(1.2), Inches(5.0), SW - Inches(2.4), Inches(0.5),
                     s["footer"], size=12, color=MUTED)
            continue

        if kind == "closing":
            add_rect(slide, 0, 0, SW, SH, WHITE)
            add_rect(slide, 0, SH - Inches(0.1), SW, Inches(0.1), BRAND)
            add_rect(slide, 0, 0, SW, Inches(0.1), SLATE)
            add_text(slide, Inches(0.5), Inches(2.8), SW - Inches(1), Inches(1.3),
                     s["title"], size=60, bold=True,
                     color=SLATE, align=PP_ALIGN.CENTER)
            add_rect(slide, (SW - Inches(2)) / 2, Inches(4.1), Inches(2), Inches(0.05), BRAND)
            add_text(slide, Inches(0.5), Inches(4.4), SW - Inches(1), Inches(0.7),
                     s["subtitle"], size=24,
                     color=SLATE_2, align=PP_ALIGN.CENTER)
            add_text(slide, Inches(0.5), Inches(6.5), SW - Inches(1), Inches(0.4),
                     s["footer"], size=11,
                     color=MUTED, align=PP_ALIGN.CENTER)
            continue

        # Common header for content slides
        header_band(slide, s["title"])
        content_top = Inches(1.1)
        content_left = Inches(0.5)
        content_w = SW - Inches(1.0)

        if kind == "bullets":
            cur_top = content_top
            if s.get("intro"):
                add_text(slide, content_left, cur_top, content_w, Inches(0.5),
                         s["intro"], size=14, color=INK)
                cur_top += Inches(0.55)
            # Bullets
            bullet_box_h = Inches(3.4)
            tb = slide.shapes.add_textbox(content_left, cur_top, content_w, bullet_box_h)
            tf = tb.text_frame
            tf.word_wrap = True
            for i, b in enumerate(s["bullets"]):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.space_after = Pt(8)
                run = p.add_run()
                run.text = f"●  {b}"
                run.font.name = "Segoe UI"
                run.font.size = Pt(16)
                run.font.color.rgb = INK
            # Highlight tiles
            if s.get("highlights"):
                tiles = s["highlights"]
                tile_w = (content_w - Inches(0.3) * (len(tiles) - 1)) / len(tiles)
                tile_h = Inches(1.1)
                tile_top = SH - Inches(1.7)
                for i, (val, label) in enumerate(tiles):
                    x = content_left + (tile_w + Inches(0.3)) * i
                    add_rect(slide, x, tile_top, tile_w, tile_h, SURFACE_2)
                    add_rect(slide, x, tile_top, Inches(0.06), tile_h, BRAND)
                    add_text(slide, x + Inches(0.2), tile_top + Inches(0.08),
                             tile_w - Inches(0.3), Inches(0.55),
                             val, size=28, bold=True, color=SLATE)
                    add_text(slide, x + Inches(0.2), tile_top + Inches(0.65),
                             tile_w - Inches(0.3), Inches(0.4),
                             label, size=11, color=MUTED)
            continue

        if kind == "table":
            cur_top = content_top
            if s.get("intro"):
                add_text(slide, content_left, cur_top, content_w, Inches(0.5),
                         s["intro"], size=13, color=MUTED)
                cur_top += Inches(0.55)
            footnote_h = Inches(0.6) if s.get("footnote") else Inches(0.1)
            tbl_h = SH - cur_top - Inches(0.4) - footnote_h
            # estimate row height
            col_widths_in = s.get("col_widths")
            build_table(slide, content_left, cur_top, content_w, tbl_h,
                        s["headers"], s["rows"], col_widths=col_widths_in,
                        font_size=10)
            if s.get("footnote"):
                add_text(slide, content_left, SH - Inches(0.95), content_w, Inches(0.55),
                         "Note: " + s["footnote"], size=10, color=MUTED)
            continue

        if kind == "layers":
            layers = s["layers"]
            n = len(layers)
            avail_h = SH - content_top - Inches(0.5)
            row_h = avail_h / n - Inches(0.12)
            for i, (tit, desc) in enumerate(layers):
                y = content_top + (row_h + Inches(0.12)) * i
                # Slate gradient from dark (top) to lighter; thin red leading bar
                shade = [SLATE, SLATE_2, SLATE_3,
                         RGBColor(0x9C, 0xA3, 0xAF), RGBColor(0xD1, 0xD5, 0xDB)]
                bar_color = shade[i % len(shade)]
                add_rect(slide, content_left, y, Inches(0.18), row_h, bar_color)
                add_rect(slide, content_left + Inches(0.18), y, content_w - Inches(0.18),
                         row_h, SURFACE)
                add_text(slide, content_left + Inches(0.4), y + Inches(0.05),
                         content_w - Inches(0.5), Inches(0.4),
                         tit, size=14, bold=True, color=SLATE)
                add_text(slide, content_left + Inches(0.4), y + Inches(0.45),
                         content_w - Inches(0.5), row_h - Inches(0.5),
                         desc, size=11, color=INK)
            continue

        if kind == "modules_p1":
            cur_top = content_top
            st = s["summary_table"]
            build_table(slide, content_left, cur_top, content_w,
                        SH - cur_top - Inches(0.95),
                        st["headers"], st["rows"], col_widths=st["col_widths"],
                        font_size=11)
            if s.get("footnote"):
                add_text(slide, content_left, SH - Inches(0.85), content_w, Inches(0.45),
                         s["footnote"], size=10, color=MUTED)
            continue

        if kind == "modules_p2":
            groups = s["groups"]
            cols = 2
            rows = (len(groups) + cols - 1) // cols
            avail_h = SH - content_top - Inches(0.5)
            avail_w = content_w
            cell_w = (avail_w - Inches(0.25)) / cols
            cell_h = (avail_h - Inches(0.2)) / rows - Inches(0.05)
            for i, (gt, desc) in enumerate(groups):
                r = i // cols
                c = i % cols
                x = content_left + (cell_w + Inches(0.25)) * c
                y = content_top + (cell_h + Inches(0.2)) * r
                add_rect(slide, x, y, Inches(0.08), cell_h, BRAND)
                add_rect(slide, x + Inches(0.08), y, cell_w - Inches(0.08), cell_h, SURFACE)
                add_text(slide, x + Inches(0.3), y + Inches(0.1), cell_w - Inches(0.4),
                         Inches(0.35), gt, size=13, bold=True, color=SLATE)
                add_text(slide, x + Inches(0.3), y + Inches(0.5), cell_w - Inches(0.4),
                         cell_h - Inches(0.55), desc, size=10, color=INK)
            continue

        if kind == "two_col":
            cur_top = content_top
            col_w = (content_w - Inches(0.3)) / 2
            col_h = SH - cur_top - Inches(0.95 if s.get("footnote") else 0.5)
            for ci, (title_key, list_key, x_offset) in enumerate([
                ("left_title", "left_bullets", 0),
                ("right_title", "right_bullets", col_w + Inches(0.3)),
            ]):
                x = content_left + x_offset
                add_rect(slide, x, cur_top, col_w, Inches(0.5), SLATE)
                add_rect(slide, x, cur_top + Inches(0.5), col_w, Inches(0.03), BRAND)
                add_text(slide, x + Inches(0.2), cur_top + Inches(0.08),
                         col_w - Inches(0.3), Inches(0.38),
                         s[title_key], size=14, bold=True,
                         color=WHITE)
                add_rect(slide, x, cur_top + Inches(0.53), col_w,
                         col_h - Inches(0.53), SURFACE)
                tb = slide.shapes.add_textbox(x + Inches(0.2), cur_top + Inches(0.6),
                                              col_w - Inches(0.4), col_h - Inches(0.7))
                tf = tb.text_frame
                tf.word_wrap = True
                for i, b in enumerate(s[list_key]):
                    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                    p.space_after = Pt(5)
                    run = p.add_run()
                    run.text = f"●  {b}"
                    run.font.name = "Segoe UI"
                    run.font.size = Pt(11)
                    run.font.color.rgb = INK
            if s.get("footnote"):
                add_text(slide, content_left, SH - Inches(0.85), content_w, Inches(0.45),
                         s["footnote"], size=10, color=MUTED)
            continue

        if kind == "diagram":
            cur_top = content_top
            if s.get("intro"):
                add_text(slide, content_left, cur_top, content_w, Inches(0.4),
                         s["intro"], size=13, color=MUTED)
                cur_top += Inches(0.45)
            # Two-col: ascii left, decisions right
            left_w = Inches(7.0)
            right_w = content_w - left_w - Inches(0.3)
            box_h = SH - cur_top - Inches(0.5)
            # ASCII box
            add_rect(slide, content_left, cur_top, left_w, box_h, SLATE)
            add_text(slide, content_left + Inches(0.2), cur_top + Inches(0.15),
                     left_w - Inches(0.3), box_h - Inches(0.2),
                     "\n".join(s["ascii"]), size=11,
                     color=WHITE, font="Consolas")
            # Decisions
            x = content_left + left_w + Inches(0.3)
            add_rect(slide, x, cur_top, right_w, Inches(0.5), SLATE)
            add_rect(slide, x, cur_top + Inches(0.5), right_w, Inches(0.03), BRAND)
            add_text(slide, x + Inches(0.2), cur_top + Inches(0.08),
                     right_w - Inches(0.3), Inches(0.38),
                     "Key Decisions", size=14, bold=True,
                     color=WHITE)
            add_rect(slide, x, cur_top + Inches(0.53), right_w, box_h - Inches(0.53), SURFACE)
            tb = slide.shapes.add_textbox(x + Inches(0.2), cur_top + Inches(0.6),
                                          right_w - Inches(0.4), box_h - Inches(0.7))
            tf = tb.text_frame
            tf.word_wrap = True
            for i, d in enumerate(s["decisions"]):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.space_after = Pt(6)
                run = p.add_run()
                run.text = f"●  {d}"
                run.font.name = "Segoe UI"
                run.font.size = Pt(11)
                run.font.color.rgb = INK
            continue

        if kind == "roadmap":
            qs = s["quarters"]
            n = len(qs)
            avail_w = content_w
            avail_h = SH - content_top - Inches(0.5)
            col_w = (avail_w - Inches(0.2) * (n - 1)) / n
            for i, (qt, items) in enumerate(qs):
                x = content_left + (col_w + Inches(0.2)) * i
                # Latest quarter gets red accent, others slate
                accent = BRAND if i == 0 else SLATE_3
                add_rect(slide, x, content_top, col_w, Inches(0.6), SLATE)
                add_rect(slide, x, content_top + Inches(0.6), col_w, Inches(0.03), accent)
                add_text(slide, x + Inches(0.15), content_top + Inches(0.13),
                         col_w - Inches(0.3), Inches(0.4),
                         qt, size=12, bold=True,
                         color=WHITE)
                add_rect(slide, x, content_top + Inches(0.63), col_w,
                         avail_h - Inches(0.63), SURFACE)
                tb = slide.shapes.add_textbox(x + Inches(0.15),
                                              content_top + Inches(0.7),
                                              col_w - Inches(0.3),
                                              avail_h - Inches(0.8))
                tf = tb.text_frame
                tf.word_wrap = True
                for j, it in enumerate(items):
                    p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
                    p.space_after = Pt(5)
                    run = p.add_run()
                    run.text = f"●  {it}"
                    run.font.name = "Segoe UI"
                    run.font.size = Pt(10)
                    run.font.color.rgb = INK
            continue

        if kind == "cta":
            cur_top = content_top
            if s.get("intro"):
                add_text(slide, content_left, cur_top, content_w, Inches(0.5),
                         s["intro"], size=14, color=MUTED)
                cur_top += Inches(0.6)
            items = s["items"]
            row_h = (SH - cur_top - Inches(0.5)) / len(items) - Inches(0.05)
            for i, (num, ttl, desc) in enumerate(items):
                y = cur_top + (row_h + Inches(0.05)) * i
                # Number circle
                circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, content_left, y + Inches(0.05),
                                              Inches(0.7), Inches(0.7))
                circ.fill.solid()
                circ.fill.fore_color.rgb = BRAND
                circ.line.fill.background()
                circ.shadow.inherit = False
                tf = circ.text_frame
                tf.margin_left = Emu(0)
                tf.margin_right = Emu(0)
                tf.margin_top = Emu(0)
                tf.margin_bottom = Emu(0)
                tf.vertical_anchor = MSO_ANCHOR.MIDDLE
                p = tf.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER
                run = p.add_run()
                run.text = num
                run.font.name = "Segoe UI"
                run.font.size = Pt(20)
                run.font.bold = True
                run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                # Title + desc
                x_text = content_left + Inches(1.0)
                add_text(slide, x_text, y + Inches(0.05),
                         content_w - Inches(1.0), Inches(0.4),
                         ttl, size=16, bold=True, color=SLATE)
                add_text(slide, x_text, y + Inches(0.45),
                         content_w - Inches(1.0), row_h - Inches(0.4),
                         desc, size=12, color=INK)
            continue

    prs.save(str(out_path))


# ======================================================================
# PDF builder (reportlab)
# ======================================================================

def build_pdf(out_path: Path):
    PAGE = landscape(A4)
    PW, PH = PAGE

    margin_l = 1.2 * cm
    margin_r = 1.2 * cm
    margin_t = 1.8 * cm   # leave room for header band
    margin_b = 1.2 * cm

    styles = getSampleStyleSheet()

    sty_title = ParagraphStyle("title", parent=styles["Title"],
                               fontName="Helvetica-Bold", fontSize=28,
                               textColor=RL_INK, leading=34, alignment=TA_LEFT)
    sty_subtitle = ParagraphStyle("subtitle", parent=styles["Normal"],
                                  fontName="Helvetica", fontSize=16,
                                  textColor=RL_SLATE_2, leading=22)
    sty_h = ParagraphStyle("h", parent=styles["Heading1"],
                           fontName="Helvetica-Bold", fontSize=18,
                           textColor=rl_colors.white, leading=22)
    sty_body = ParagraphStyle("body", parent=styles["Normal"],
                              fontName="Helvetica", fontSize=11,
                              textColor=RL_INK, leading=15)
    sty_muted = ParagraphStyle("muted", parent=styles["Normal"],
                               fontName="Helvetica-Oblique", fontSize=9,
                               textColor=RL_MUTED, leading=12)
    sty_bullet = ParagraphStyle("bullet", parent=sty_body, leftIndent=14,
                                bulletIndent=2, fontSize=11, leading=15)
    sty_col_h = ParagraphStyle("colh", parent=styles["Normal"],
                               fontName="Helvetica-Bold", fontSize=12,
                               textColor=rl_colors.white, leading=14)
    sty_small = ParagraphStyle("small", parent=sty_body, fontSize=9, leading=11)
    sty_tile_val = ParagraphStyle("tileval", parent=styles["Normal"],
                                  fontName="Helvetica-Bold", fontSize=22,
                                  textColor=RL_SLATE, alignment=TA_CENTER, leading=26)
    sty_tile_lbl = ParagraphStyle("tilelbl", parent=styles["Normal"],
                                  fontName="Helvetica", fontSize=9,
                                  textColor=RL_INK, alignment=TA_CENTER, leading=11)
    sty_quarter_h = ParagraphStyle("qh", parent=styles["Normal"],
                                   fontName="Helvetica-Bold", fontSize=11,
                                   textColor=rl_colors.white, leading=13)
    sty_layer_t = ParagraphStyle("lt", parent=styles["Normal"],
                                 fontName="Helvetica-Bold", fontSize=12,
                                 textColor=RL_SLATE, leading=14)
    sty_layer_d = ParagraphStyle("ld", parent=styles["Normal"],
                                 fontName="Helvetica", fontSize=9.5,
                                 textColor=RL_INK, leading=12)
    sty_code = ParagraphStyle("code", parent=styles["Code"],
                              fontName="Courier", fontSize=9,
                              textColor=rl_colors.HexColor("#FAFAF9"), leading=12)

    def page_decoration(canvas, doc):
        # Top slate band with thin red accent line
        canvas.saveState()
        canvas.setFillColor(RL_SLATE)
        canvas.rect(0, PH - 1.2 * cm, PW, 1.2 * cm, fill=1, stroke=0)
        canvas.setFillColor(RL_BRAND)
        canvas.rect(0, PH - 1.25 * cm, PW, 0.08 * cm, fill=1, stroke=0)
        # Footer band
        canvas.setFillColor(RL_SURFACE_2)
        canvas.rect(0, 0, PW, 0.6 * cm, fill=1, stroke=0)
        canvas.setFillColor(RL_MUTED)
        canvas.setFont("Helvetica", 8)
        canvas.drawString(margin_l, 0.2 * cm,
                          "Erajaya VAS  ·  Odoo 19 Platform  ·  Mei 2026")
        canvas.drawRightString(PW - margin_r, 0.2 * cm,
                               f"Confidential — Internal  ·  {doc.page}")
        canvas.restoreState()

    def title_page_decoration(canvas, doc):
        canvas.saveState()
        canvas.setFillColor(rl_colors.white)
        canvas.rect(0, 0, PW, PH, fill=1, stroke=0)
        canvas.setFillColor(RL_SLATE)
        canvas.rect(0, 0, 0.7 * cm, PH, fill=1, stroke=0)
        canvas.setFillColor(RL_BRAND)
        canvas.rect(0.7 * cm, 0, 0.15 * cm, PH, fill=1, stroke=0)
        canvas.restoreState()

    flowables = []

    def page_header_para(title):
        # The header band is drawn by page_decoration; we add a Paragraph that
        # sits inside it via a Spacer + Paragraph trick: we use a Table.
        t = Table([[Paragraph(title, sty_h)]],
                  colWidths=[PW - margin_l - margin_r],
                  rowHeights=[0.9 * cm])
        t.setStyle(TableStyle([
            ("BACKGROUND", (0, 0), (-1, -1), RL_SLATE),
            ("LINEBELOW", (0, 0), (-1, -1), 1.5, RL_BRAND),
            ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
            ("LEFTPADDING", (0, 0), (-1, -1), 6),
            ("RIGHTPADDING", (0, 0), (-1, -1), 6),
            ("TOPPADDING", (0, 0), (-1, -1), 4),
            ("BOTTOMPADDING", (0, 0), (-1, -1), 4),
        ]))
        return t

    def std_table(headers, rows, col_widths=None, font_size=9):
        n_cols = len(headers)
        avail = PW - margin_l - margin_r
        if col_widths:
            total = sum(col_widths)
            widths = [avail * w / total for w in col_widths]
        else:
            widths = [avail / n_cols] * n_cols
        h_paras = [Paragraph(f"<b>{h}</b>", ParagraphStyle(
            "th", fontName="Helvetica-Bold", fontSize=font_size + 1,
            textColor=rl_colors.white, leading=font_size + 3)) for h in headers]
        body = [h_paras]
        for r in rows:
            body.append([Paragraph(str(c), ParagraphStyle(
                "td", fontName="Helvetica", fontSize=font_size,
                textColor=RL_INK, leading=font_size + 2)) for c in r])
        tbl = Table(body, colWidths=widths, repeatRows=1)
        ts = [
            ("BACKGROUND", (0, 0), (-1, 0), RL_SLATE),
            ("VALIGN", (0, 0), (-1, -1), "TOP"),
            ("GRID", (0, 0), (-1, -1), 0.4, RL_BORDER),
            ("LEFTPADDING", (0, 0), (-1, -1), 5),
            ("RIGHTPADDING", (0, 0), (-1, -1), 5),
            ("TOPPADDING", (0, 0), (-1, -1), 4),
            ("BOTTOMPADDING", (0, 0), (-1, -1), 4),
        ]
        # alt row shading
        for i in range(1, len(body)):
            if i % 2 == 1:
                ts.append(("BACKGROUND", (0, i), (-1, i), RL_SURFACE_2))
            else:
                ts.append(("BACKGROUND", (0, i), (-1, i), rl_colors.white))
        tbl.setStyle(TableStyle(ts))
        return tbl

    # We will manually manage one slide per page by using PageBreak.
    # First page = title with custom canvas; remaining pages use page_decoration.
    # ReportLab limitation: onFirstPage / onLaterPages — we need different handler.

    title_slide = SLIDES[0]
    rest = SLIDES[1:]

    avail_w = PW - margin_l - margin_r
    avail_h = PH - margin_t - margin_b

    # ---- Title page (handled via flowables on the white canvas) ----
    flowables.append(Spacer(1, 5.5 * cm))
    flowables.append(Paragraph(title_slide["title"], sty_title))
    flowables.append(Spacer(1, 0.4 * cm))
    flowables.append(Paragraph(title_slide["subtitle"], sty_subtitle))
    flowables.append(Spacer(1, 0.5 * cm))
    # Horizontal accent
    accent = Table([[""]], colWidths=[6 * cm], rowHeights=[0.12 * cm])
    accent.setStyle(TableStyle([("BACKGROUND", (0, 0), (-1, -1), RL_BRAND)]))
    flowables.append(accent)
    flowables.append(Spacer(1, 0.4 * cm))
    flowables.append(Paragraph(title_slide["footer"], sty_muted))
    flowables.append(PageBreak())

    # ---- Content slides ----
    for s in rest:
        kind = s["kind"]
        flowables.append(page_header_para(s["title"]))
        flowables.append(Spacer(1, 0.4 * cm))

        if kind == "bullets":
            if s.get("intro"):
                flowables.append(Paragraph(s["intro"], sty_body))
                flowables.append(Spacer(1, 0.3 * cm))
            for b in s["bullets"]:
                flowables.append(Paragraph(f"●&nbsp;&nbsp;{b}", sty_bullet))
                flowables.append(Spacer(1, 0.12 * cm))
            if s.get("highlights"):
                flowables.append(Spacer(1, 0.4 * cm))
                tiles = s["highlights"]
                tile_data = [[Paragraph(v, sty_tile_val) for v, _ in tiles],
                             [Paragraph(l, sty_tile_lbl) for _, l in tiles]]
                w_each = (avail_w - 0.6 * cm) / len(tiles)
                ttbl = Table(tile_data,
                             colWidths=[w_each] * len(tiles),
                             rowHeights=[1.3 * cm, 0.7 * cm])
                ts = [
                    ("BACKGROUND", (0, 0), (-1, -1), RL_SURFACE_2),
                    ("LINEABOVE", (0, 0), (-1, 0), 2, RL_BRAND),
                    ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
                    ("ALIGN", (0, 0), (-1, -1), "CENTER"),
                    ("TOPPADDING", (0, 0), (-1, -1), 4),
                    ("BOTTOMPADDING", (0, 0), (-1, -1), 4),
                ]
                ttbl.setStyle(TableStyle(ts))
                flowables.append(ttbl)

        elif kind == "table":
            if s.get("intro"):
                flowables.append(Paragraph(s["intro"], sty_muted))
                flowables.append(Spacer(1, 0.25 * cm))
            flowables.append(std_table(s["headers"], s["rows"],
                                       col_widths=s.get("col_widths"),
                                       font_size=9))
            if s.get("footnote"):
                flowables.append(Spacer(1, 0.25 * cm))
                flowables.append(Paragraph("<i>Note:</i> " + s["footnote"], sty_muted))

        elif kind == "layers":
            n = len(s["layers"])
            row_h = (avail_h - 1.0 * cm) / n
            for i, (tit, desc) in enumerate(s["layers"]):
                bar_color = [RL_SLATE, RL_SLATE_2, RL_SLATE_3,
                             rl_colors.HexColor("#9CA3AF"),
                             rl_colors.HexColor("#D1D5DB")][i % 5]
                tbl = Table([[
                    "", Paragraph(f"<b>{tit}</b><br/>"
                                  f"<font size=9 color='#111827'>{desc}</font>",
                                  sty_layer_t)
                ]], colWidths=[0.4 * cm, avail_w - 0.4 * cm],
                    rowHeights=[row_h - 0.18 * cm])
                tbl.setStyle(TableStyle([
                    ("BACKGROUND", (0, 0), (0, 0), bar_color),
                    ("BACKGROUND", (1, 0), (1, 0), RL_SURFACE),
                    ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
                    ("LEFTPADDING", (1, 0), (1, 0), 10),
                    ("RIGHTPADDING", (1, 0), (1, 0), 8),
                    ("TOPPADDING", (0, 0), (-1, -1), 6),
                    ("BOTTOMPADDING", (0, 0), (-1, -1), 6),
                ]))
                flowables.append(tbl)
                flowables.append(Spacer(1, 0.12 * cm))

        elif kind == "modules_p1":
            st = s["summary_table"]
            flowables.append(std_table(st["headers"], st["rows"],
                                       col_widths=st["col_widths"],
                                       font_size=10))
            if s.get("footnote"):
                flowables.append(Spacer(1, 0.25 * cm))
                flowables.append(Paragraph(s["footnote"], sty_muted))

        elif kind == "modules_p2":
            # 2-col grid of cards
            groups = s["groups"]
            cards = []
            for gt, desc in groups:
                inner = Table([[Paragraph(f"<b><font color='#1F2937'>{gt}</font></b><br/>"
                                          f"<font size=9>{desc}</font>",
                                          sty_layer_t)]],
                              colWidths=[(avail_w - 0.6 * cm) / 2 - 0.4 * cm])
                inner.setStyle(TableStyle([
                    ("LEFTPADDING", (0, 0), (-1, -1), 8),
                    ("RIGHTPADDING", (0, 0), (-1, -1), 8),
                    ("TOPPADDING", (0, 0), (-1, -1), 6),
                    ("BOTTOMPADDING", (0, 0), (-1, -1), 6),
                ]))
                cards.append(inner)
            # build pairs
            data = []
            for i in range(0, len(cards), 2):
                row = cards[i:i+2]
                if len(row) < 2:
                    row.append("")
                data.append(row)
            col_w = (avail_w - 0.6 * cm) / 2
            grid = Table(data, colWidths=[col_w, col_w])
            ts = [
                ("VALIGN", (0, 0), (-1, -1), "TOP"),
                ("LEFTPADDING", (0, 0), (-1, -1), 0),
                ("RIGHTPADDING", (0, 0), (-1, -1), 6),
                ("TOPPADDING", (0, 0), (-1, -1), 0),
                ("BOTTOMPADDING", (0, 0), (-1, -1), 6),
            ]
            # Each card gets a left brand bar via background trick is hard; use BOX
            for r in range(len(data)):
                for c in range(2):
                    ts.append(("LINEBEFORE", (c, r), (c, r), 3, RL_BRAND))
                    ts.append(("BACKGROUND", (c, r), (c, r), RL_SURFACE))
            grid.setStyle(TableStyle(ts))
            flowables.append(grid)

        elif kind == "two_col":
            col_w = (avail_w - 0.4 * cm) / 2
            def col_block(title, bullets):
                inner_rows = [[Paragraph(title, sty_col_h)]]
                bullet_paras = []
                for b in bullets:
                    bullet_paras.append(Paragraph(f"●&nbsp;&nbsp;{b}", sty_bullet))
                    bullet_paras.append(Spacer(1, 0.08 * cm))
                inner_rows.append([bullet_paras])
                t = Table(inner_rows, colWidths=[col_w])
                t.setStyle(TableStyle([
                    ("BACKGROUND", (0, 0), (-1, 0), RL_SLATE),
                    ("LINEBELOW", (0, 0), (-1, 0), 1.5, RL_BRAND),
                    ("BACKGROUND", (0, 1), (-1, 1), RL_SURFACE),
                    ("LEFTPADDING", (0, 0), (-1, -1), 8),
                    ("RIGHTPADDING", (0, 0), (-1, -1), 8),
                    ("TOPPADDING", (0, 0), (-1, -1), 6),
                    ("BOTTOMPADDING", (0, 0), (-1, -1), 6),
                    ("VALIGN", (0, 1), (-1, 1), "TOP"),
                ]))
                return t
            grid = Table([[
                col_block(s["left_title"], s["left_bullets"]),
                col_block(s["right_title"], s["right_bullets"]),
            ]], colWidths=[col_w, col_w])
            grid.setStyle(TableStyle([
                ("VALIGN", (0, 0), (-1, -1), "TOP"),
                ("LEFTPADDING", (0, 0), (-1, -1), 0),
                ("RIGHTPADDING", (0, 0), (-1, -1), 6),
                ("TOPPADDING", (0, 0), (-1, -1), 0),
            ]))
            flowables.append(grid)
            if s.get("footnote"):
                flowables.append(Spacer(1, 0.3 * cm))
                flowables.append(Paragraph(s["footnote"], sty_muted))

        elif kind == "diagram":
            if s.get("intro"):
                flowables.append(Paragraph(s["intro"], sty_muted))
                flowables.append(Spacer(1, 0.25 * cm))
            left_w = avail_w * 0.6
            right_w = avail_w - left_w - 0.4 * cm
            ascii_text = "<br/>".join(
                line.replace(" ", "&nbsp;") for line in s["ascii"])
            ascii_para = Paragraph(ascii_text, sty_code)
            left_cell = Table([[ascii_para]], colWidths=[left_w])
            left_cell.setStyle(TableStyle([
                ("BACKGROUND", (0, 0), (-1, -1), RL_SLATE),
                ("LEFTPADDING", (0, 0), (-1, -1), 10),
                ("RIGHTPADDING", (0, 0), (-1, -1), 10),
                ("TOPPADDING", (0, 0), (-1, -1), 10),
                ("BOTTOMPADDING", (0, 0), (-1, -1), 10),
            ]))
            dec_inner = [[Paragraph("Key Decisions", sty_col_h)]]
            dec_bullets = []
            for d in s["decisions"]:
                dec_bullets.append(Paragraph(f"●&nbsp;&nbsp;{d}", sty_bullet))
                dec_bullets.append(Spacer(1, 0.1 * cm))
            dec_inner.append([dec_bullets])
            right_cell = Table(dec_inner, colWidths=[right_w])
            right_cell.setStyle(TableStyle([
                ("BACKGROUND", (0, 0), (-1, 0), RL_SLATE),
                ("LINEBELOW", (0, 0), (-1, 0), 1.5, RL_BRAND),
                ("BACKGROUND", (0, 1), (-1, 1), RL_SURFACE),
                ("LEFTPADDING", (0, 0), (-1, -1), 8),
                ("RIGHTPADDING", (0, 0), (-1, -1), 8),
                ("TOPPADDING", (0, 0), (-1, -1), 6),
                ("BOTTOMPADDING", (0, 0), (-1, -1), 6),
                ("VALIGN", (0, 1), (-1, 1), "TOP"),
            ]))
            grid = Table([[left_cell, right_cell]],
                         colWidths=[left_w, right_w])
            grid.setStyle(TableStyle([
                ("VALIGN", (0, 0), (-1, -1), "TOP"),
                ("LEFTPADDING", (0, 0), (-1, -1), 0),
                ("RIGHTPADDING", (0, 0), (-1, -1), 6),
            ]))
            flowables.append(grid)

        elif kind == "roadmap":
            qs = s["quarters"]
            n = len(qs)
            col_w = (avail_w - 0.4 * cm * (n - 1)) / n
            cells = []
            for qi, (qt, items) in enumerate(qs):
                inner = [[Paragraph(qt, sty_quarter_h)]]
                bps = []
                for it in items:
                    bps.append(Paragraph(f"●&nbsp;&nbsp;{it}", sty_small))
                    bps.append(Spacer(1, 0.08 * cm))
                inner.append([bps])
                t = Table(inner, colWidths=[col_w])
                accent = RL_BRAND if qi == 0 else RL_SLATE_3
                t.setStyle(TableStyle([
                    ("BACKGROUND", (0, 0), (-1, 0), RL_SLATE),
                    ("LINEBELOW", (0, 0), (-1, 0), 1.5, accent),
                    ("BACKGROUND", (0, 1), (-1, 1), RL_SURFACE),
                    ("LEFTPADDING", (0, 0), (-1, -1), 6),
                    ("RIGHTPADDING", (0, 0), (-1, -1), 6),
                    ("TOPPADDING", (0, 0), (-1, -1), 5),
                    ("BOTTOMPADDING", (0, 0), (-1, -1), 5),
                    ("VALIGN", (0, 1), (-1, 1), "TOP"),
                ]))
                cells.append(t)
            grid = Table([cells], colWidths=[col_w] * n)
            grid.setStyle(TableStyle([
                ("VALIGN", (0, 0), (-1, -1), "TOP"),
                ("LEFTPADDING", (0, 0), (-1, -1), 0),
                ("RIGHTPADDING", (0, 0), (-1, -1), 6),
            ]))
            flowables.append(grid)

        elif kind == "cta":
            if s.get("intro"):
                flowables.append(Paragraph(s["intro"], sty_muted))
                flowables.append(Spacer(1, 0.3 * cm))
            for num, ttl, desc in s["items"]:
                num_cell = Paragraph(
                    f"<para align='center'><font color='white' size=16><b>{num}</b></font></para>",
                    sty_body)
                txt_cell = Paragraph(
                    f"<font color='#1F2937' size=13><b>{ttl}</b></font><br/>"
                    f"<font size=10>{desc}</font>", sty_body)
                row = Table([[num_cell, txt_cell]],
                            colWidths=[1.5 * cm, avail_w - 1.5 * cm])
                row.setStyle(TableStyle([
                    ("BACKGROUND", (0, 0), (0, 0), RL_SLATE),
                    ("LINEAFTER", (0, 0), (0, 0), 2, RL_BRAND),
                    ("BACKGROUND", (1, 0), (1, 0), RL_SURFACE),
                    ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
                    ("LEFTPADDING", (1, 0), (1, 0), 10),
                    ("RIGHTPADDING", (1, 0), (1, 0), 8),
                    ("TOPPADDING", (0, 0), (-1, -1), 8),
                    ("BOTTOMPADDING", (0, 0), (-1, -1), 8),
                ]))
                flowables.append(row)
                flowables.append(Spacer(1, 0.2 * cm))

        elif kind == "closing":
            flowables.append(Spacer(1, 4 * cm))
            flowables.append(Paragraph(
                f"<para align='center'><font size=40 color='#1F2937'><b>{s['title']}</b></font></para>",
                sty_body))
            flowables.append(Spacer(1, 0.4 * cm))
            accent = Table([[""]], colWidths=[4 * cm], rowHeights=[0.1 * cm])
            accent.setStyle(TableStyle([
                ("BACKGROUND", (0, 0), (-1, -1), RL_BRAND),
                ("ALIGN", (0, 0), (-1, -1), "CENTER"),
            ]))
            inner_wrap = Table([[accent]], colWidths=[avail_w])
            inner_wrap.setStyle(TableStyle([
                ("ALIGN", (0, 0), (-1, -1), "CENTER"),
                ("LEFTPADDING", (0, 0), (-1, -1), 0),
                ("RIGHTPADDING", (0, 0), (-1, -1), 0),
            ]))
            flowables.append(inner_wrap)
            flowables.append(Spacer(1, 0.5 * cm))
            flowables.append(Paragraph(
                f"<para align='center'><font size=18 color='#374151'>{s['subtitle']}</font></para>",
                sty_body))
            flowables.append(Spacer(1, 2 * cm))
            flowables.append(Paragraph(
                f"<para align='center'><font size=10 color='#6B7280'>{s['footer']}</font></para>",
                sty_body))

        flowables.append(PageBreak())

    # remove final extra page break
    if flowables and isinstance(flowables[-1], PageBreak):
        flowables.pop()

    # Custom handler — first page is title, rest get content decoration.
    def on_first(canvas, doc):
        title_page_decoration(canvas, doc)
    def on_later(canvas, doc):
        page_decoration(canvas, doc)

    doc = SimpleDocTemplate(str(out_path), pagesize=PAGE,
                            leftMargin=margin_l, rightMargin=margin_r,
                            topMargin=margin_t, bottomMargin=margin_b,
                            title="Erajaya VAS — Odoo Platform")
    doc.build(flowables, onFirstPage=on_first, onLaterPages=on_later)


def main():
    out_dir = Path(__file__).resolve().parent.parent / "docs"
    out_dir.mkdir(exist_ok=True)
    pptx_path = out_dir / "presentation-erajaya-vas.pptx"
    pdf_path = out_dir / "presentation-erajaya-vas.pdf"

    print(f"Building PPTX -> {pptx_path}")
    build_pptx(pptx_path)
    print(f"  OK {pptx_path.stat().st_size:,} bytes")

    print(f"Building PDF  -> {pdf_path}")
    build_pdf(pdf_path)
    print(f"  OK {pdf_path.stat().st_size:,} bytes")


if __name__ == "__main__":
    main()
